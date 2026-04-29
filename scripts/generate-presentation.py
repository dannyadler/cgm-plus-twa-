#!/usr/bin/env python3
"""
Generate the BioT-branded PowerPoint deck for the AI-Assisted Regulatory
Documentation customer presentation.

Output: presentations/trinity-cgm-plus-ai-regulatory.pptx

Brand reference: bioT BrandBook 2024.
- Primary: Purple 100 #4545D8
- Dark variant: Black #141437 (deep navy-purple)
- Accent: Peach 100 #FFA5A0
- Backgrounds: Grey 100 #EDEDED, White #FFFFFF
- Typography: Helvetica Neue (headings/titles), DM Sans (body/quotes)

Layout idiom (mirrors the BrandBook slides):
- 16:9 widescreen.
- Section dividers: full Purple 100 background, peach circle with section
  number top-left of title block, big white title bottom-left, vertical date
  text rotated on right edge.
- Content slides: Grey 100 background, top-left title with horizontal divider
  line extending to the right, bottom "bioT" wordmark with divider line.
- Closing slide: Black/navy background, big white "Thank You." bottom-left.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ── Brand constants ──────────────────────────────────────────────────────────

PURPLE_100 = RGBColor(0x45, 0x45, 0xD8)   # primary
PURPLE_120 = RGBColor(0x3D, 0x3E, 0xD1)
PURPLE_75  = RGBColor(0x6F, 0x6F, 0xFA)
PURPLE_50  = RGBColor(0xC0, 0xC0, 0xEE)
PURPLE_15  = RGBColor(0xDB, 0xDB, 0xE6)
PEACH_100  = RGBColor(0xFF, 0xA5, 0xA0)
PEACH_50   = RGBColor(0xFF, 0xCB, 0xC8)
BLACK_NAVY = RGBColor(0x14, 0x14, 0x37)   # deep navy used as "Black"
GREY_200   = RGBColor(0xE0, 0xE0, 0xE0)
GREY_100   = RGBColor(0xED, 0xED, 0xED)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
INK        = RGBColor(0x1A, 0x1A, 0x1A)
MUTED      = RGBColor(0x66, 0x66, 0x66)

FONT_HEAD  = "Helvetica Neue"
FONT_BODY  = "DM Sans"

# 16:9 page in EMU (914400 EMU/inch)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FOOTER_TEXT = "Trinity CGM Plus  |  AI-Assisted Regulatory Documentation  |  2026"

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, rgb):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill_rgb, line=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    if not line:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_line(slide, x1, y1, x2, y2, rgb, weight_pt=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2 - x1, y2 - y1)  # straight
    line.line.color.rgb = rgb
    line.line.width = Pt(weight_pt)
    return line


def add_text(slide, x, y, w, h, text, *, font=FONT_HEAD, size=18, bold=False,
             italic=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = ln
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_circle(slide, x, y, d, fill_rgb, label=None, label_color=INK,
               label_size=14, label_bold=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_rgb
    sh.line.fill.background()
    if label is not None:
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = FONT_HEAD
        run.font.size = Pt(label_size)
        run.font.bold = label_bold
        run.font.color.rgb = label_color
    return sh


def add_logo_mark(slide, x, y, size, color=WHITE):
    """Recreate a simplified bioT mark (rounded square + plus). Replace with
    the official asset if available."""
    # purple rounded square as background swatch (when on grey)
    if color == INK or color == BLACK_NAVY:
        sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
        sw.adjustments[0] = 0.20
        sw.fill.solid()
        sw.fill.fore_color.rgb = PURPLE_100
        sw.line.fill.background()
        sw.shadow.inherit = False
        plus_color = WHITE
    else:
        # When mark sits on a purple/dark slide, just draw a white square outline + plus
        sw = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size)
        sw.adjustments[0] = 0.20
        sw.fill.solid()
        sw.fill.fore_color.rgb = WHITE
        sw.line.fill.background()
        sw.shadow.inherit = False
        plus_color = PURPLE_100

    # Plus shape inside
    p_w = int(size * 0.55)
    p_t = int(size * 0.16)
    cx = x + size // 2
    cy = y + size // 2
    h_arm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   cx - p_w // 2, cy - p_t // 2, p_w, p_t)
    v_arm = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   cx - p_t // 2, cy - p_w // 2, p_t, p_w)
    for arm in (h_arm, v_arm):
        arm.fill.solid()
        arm.fill.fore_color.rgb = plus_color
        arm.line.fill.background()
        arm.shadow.inherit = False

    # wordmark "bioT" text next to mark
    add_text(slide, x + int(size * 1.15), y, Inches(2), size,
             "bioT", font=FONT_HEAD, size=int(size / Emu(1) / 12000) or 28,
             bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, on_dark=False):
    """Bottom-left mini wordmark + horizontal divider line + bottom-right
    page meta (mirrors BrandBook's 'bioT BrandBook | 2024')."""
    color = WHITE if on_dark else INK
    line_color = WHITE if on_dark else INK
    # mini logo
    add_logo_mark(slide, Inches(0.55), Inches(7.0), Inches(0.32), color=color)
    # divider line
    add_line(slide, Inches(2.2), Inches(7.16), Inches(12.6), Inches(7.16),
             line_color, weight_pt=0.75)


def add_top_divider(slide, title_text, on_dark=False):
    """Top-left section caption with horizontal divider extending right
    (mirrors BrandBook content slides)."""
    color = WHITE if on_dark else INK
    line_color = WHITE if on_dark else INK
    add_text(slide, Inches(0.55), Inches(0.4), Inches(3.5), Inches(0.4),
             title_text, font=FONT_HEAD, size=16, bold=True, color=color)
    add_line(slide, Inches(4.3), Inches(0.6), Inches(12.6), Inches(0.6),
             line_color, weight_pt=0.75)


def add_right_meta(slide, on_dark=False):
    """Vertical meta on right edge. python-pptx rotation requires raw XML."""
    color = WHITE if on_dark else INK
    tb = add_text(slide, Inches(12.85), Inches(2.2), Inches(0.4), Inches(3.5),
                  FOOTER_TEXT, font=FONT_HEAD, size=8, color=color)
    # rotate 270deg
    sp = tb._element
    spPr = sp.find('.//' + qn('p:spPr'))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn('a:xfrm'))
    xfrm.set('rot', '-5400000')  # -90 deg in 60000-th of a degree


# ── Slide builders ───────────────────────────────────────────────────────────

def build_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s, PURPLE_100)
    # mark
    add_logo_mark(s, Inches(0.6), Inches(0.6), Inches(0.55), color=WHITE)
    # peach circle marker
    add_circle(s, Inches(0.55), Inches(3.5), Inches(0.85),
               PEACH_100, label="00", label_color=BLACK_NAVY, label_size=16)
    # big title
    add_text(s, Inches(0.55), Inches(4.6), Inches(11), Inches(2),
             ["AI-Assisted",
              "Regulatory Documentation"],
             font=FONT_HEAD, size=64, bold=True, color=WHITE,
             line_spacing=0.95)
    # subtitle
    add_text(s, Inches(0.55), Inches(6.55), Inches(11), Inches(0.5),
             "Trinity CGM Plus  ·  Presale demonstration  ·  April 2026",
             font=FONT_BODY, size=14, color=PURPLE_15)
    add_right_meta(s, on_dark=True)
    return s


def build_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    # left purple block with "What's in"
    add_rect(s, Inches(0), Inches(0), Inches(4.5), SLIDE_H, PURPLE_100)
    add_text(s, Inches(0.55), Inches(5.6), Inches(4), Inches(1.6),
             "What's in", font=FONT_HEAD, size=56, bold=True, color=PURPLE_15,
             line_spacing=0.95)
    # right side numbered list
    items = [
        ("01.", "The product"),
        ("02.", "DHF artifact set (9 documents)"),
        ("03.", "The CI pipeline"),
        ("04.", "Where AI ends, humans begin"),
        ("05.", "Quantified value"),
        ("06.", "Q&A"),
    ]
    y = Inches(1.25)
    for num, label in items:
        add_text(s, Inches(5.2), y, Inches(1.1), Inches(0.6),
                 num, font=FONT_HEAD, size=22, bold=True, color=INK)
        add_text(s, Inches(6.25), y, Inches(7), Inches(0.6),
                 label, font=FONT_HEAD, size=22, color=INK)
        y += Inches(0.85)
    add_footer(s)
    add_right_meta(s)
    return s


def build_section_divider(prs, num, title_lines):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, PURPLE_100)
    add_circle(s, Inches(0.55), Inches(4.55), Inches(0.85),
               PEACH_100, label=num, label_color=BLACK_NAVY, label_size=16)
    add_text(s, Inches(0.55), Inches(5.5), Inches(11), Inches(2),
             title_lines, font=FONT_HEAD, size=72, bold=True, color=PURPLE_15,
             line_spacing=0.95)
    add_right_meta(s, on_dark=True)
    return s


def build_content(prs, top_caption, title, body_blocks):
    """Generic content slide. body_blocks = list of dicts with x,y,w,h,text,opts."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    add_top_divider(s, top_caption)
    # title
    add_text(s, Inches(0.55), Inches(1.0), Inches(12), Inches(1),
             title, font=FONT_HEAD, size=42, bold=True, color=INK,
             line_spacing=1.0)
    for blk in body_blocks:
        add_text(s, blk["x"], blk["y"], blk["w"], blk["h"], blk["text"],
                 font=blk.get("font", FONT_BODY),
                 size=blk.get("size", 14),
                 bold=blk.get("bold", False),
                 italic=blk.get("italic", False),
                 color=blk.get("color", INK),
                 align=blk.get("align", PP_ALIGN.LEFT),
                 line_spacing=blk.get("line_spacing", 1.25))
    add_footer(s)
    add_right_meta(s)
    return s


def build_thank_you(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BLACK_NAVY)
    add_logo_mark(s, Inches(0.55), Inches(0.5), Inches(0.55), color=WHITE)
    add_text(s, Inches(0.55), Inches(5.6), Inches(11), Inches(1.6),
             "Thank You.", font=FONT_HEAD, size=72, bold=True, color=PURPLE_15,
             line_spacing=0.95)
    add_text(s, Inches(0.55), Inches(6.7), Inches(11), Inches(0.4),
             "daniel@biot-med.com  ·  github.com/dannyadler/cgm-plus-twa-",
             font=FONT_BODY, size=12, color=PURPLE_15)
    add_right_meta(s, on_dark=True)
    return s


# ── Specific content slides ──────────────────────────────────────────────────

def build_product_context(prs):
    s = build_content(
        prs,
        "Section 01  ·  The product",
        "Trinity CGM Plus presale PWA",
        [
            {"x": Inches(0.55), "y": Inches(2.2), "w": Inches(12), "h": Inches(0.6),
             "text": "What it is", "size": 13, "bold": True, "color": MUTED},
            {"x": Inches(0.55), "y": Inches(2.6), "w": Inches(12), "h": Inches(1.3),
             "text": ("Single-file vanilla-JS Progressive Web App on top of BioT Demo2. "
                      "Deployed to Netlify, wrapped as Android Trusted Web Activity. "
                      "Five tabs displaying glucose, HR, activity, skin temp from BioT."),
             "size": 16},

            {"x": Inches(0.55), "y": Inches(4.0), "w": Inches(12), "h": Inches(0.6),
             "text": "What changed in the last week", "size": 13, "bold": True, "color": MUTED},
            {"x": Inches(0.55), "y": Inches(4.4), "w": Inches(12), "h": Inches(2.0),
             "text": ("• Replaced hardcoded credentials in client JS with a Netlify Function login proxy.\n"
                      "• Added WebAuthn fingerprint unlock backed by the platform authenticator.\n"
                      "• Generated nine regulatory documents and a real CI pipeline — what we'll walk through."),
             "size": 15, "line_spacing": 1.4},
        ],
    )
    return s


def build_dhf_overview(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    add_top_divider(s, "Section 02  ·  DHF artifact set")
    add_text(s, Inches(0.55), Inches(1.0), Inches(12), Inches(0.8),
             "Nine documents, ~30 minutes of AI generation",
             font=FONT_HEAD, size=34, bold=True, color=INK)
    rows = [
        ("01", "Software Safety Classification", "IEC 62304 §4.3", "1–2 d"),
        ("02", "Software Requirements Specification", "IEC 62304 §5.2", "3–5 d"),
        ("03", "Software Design Description", "IEC 62304 §5.3", "3–5 d"),
        ("04", "Risk Analysis", "ISO 14971:2019", "5–10 d"),
        ("05", "Software Test Description", "IEC 62304 §5.5–5.7", "5–10 d"),
        ("06", "Cybersecurity Assessment", "FDA Premarket 2023", "4–7 d"),
        ("07", "SBOM", "CycloneDX 1.5", "0.5 d"),
        ("08", "Traceability Matrix", "IEC 62304 implicit", "2–4 d"),
        ("09", "Talk Track + Slides", "—", "—"),
    ]
    # header
    headers = [("#", Inches(0.55)), ("Document", Inches(1.05)),
               ("Standard", Inches(7.4)), ("Manual effort", Inches(11.0))]
    for h, x in headers:
        add_text(s, x, Inches(2.1), Inches(3.5), Inches(0.4), h,
                 font=FONT_HEAD, size=11, bold=True, color=MUTED)
    # rows
    y = Inches(2.55)
    for num, doc, std, eff in rows:
        add_text(s, Inches(0.55), y, Inches(0.5), Inches(0.4), num,
                 font=FONT_HEAD, size=13, bold=True, color=PURPLE_100)
        add_text(s, Inches(1.05), y, Inches(6.3), Inches(0.4), doc,
                 font=FONT_HEAD, size=13, color=INK)
        add_text(s, Inches(7.4), y, Inches(3.5), Inches(0.4), std,
                 font=FONT_BODY, size=12, color=MUTED)
        add_text(s, Inches(11.0), y, Inches(2), Inches(0.4), eff,
                 font=FONT_BODY, size=12, color=INK)
        y += Inches(0.45)
    add_text(s, Inches(0.55), Inches(6.7), Inches(12), Inches(0.4),
             "Equivalent manual effort across all nine: 25–45 person-days plus clinical and security review.",
             font=FONT_BODY, size=11, italic=True, color=MUTED)
    add_footer(s)
    add_right_meta(s)
    return s


def build_doc_card(prs, num, full_title, standard, key_points, money_quote):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    add_top_divider(s, f"Section 02.{num}  ·  TCG-REG-{num}")
    # title row with peach disc
    add_circle(s, Inches(0.55), Inches(1.0), Inches(0.6), PEACH_100,
               label=num, label_color=BLACK_NAVY, label_size=14)
    add_text(s, Inches(1.35), Inches(1.0), Inches(11.5), Inches(0.7),
             full_title, font=FONT_HEAD, size=30, bold=True, color=INK)
    add_text(s, Inches(1.35), Inches(1.7), Inches(11.5), Inches(0.4),
             standard, font=FONT_BODY, size=12, italic=True, color=MUTED)
    # key points list
    body = ""
    for pt in key_points:
        body += "•  " + pt + "\n"
    add_text(s, Inches(0.55), Inches(2.6), Inches(8.0), Inches(4.0),
             body.rstrip(), font=FONT_BODY, size=14, line_spacing=1.45, color=INK)
    # money-quote box on right
    add_rect(s, Inches(8.9), Inches(2.6), Inches(3.9), Inches(3.5), WHITE)
    add_text(s, Inches(9.1), Inches(2.7), Inches(3.5), Inches(0.4),
             "What this slide says", font=FONT_HEAD, size=10, bold=True,
             color=PURPLE_100)
    add_text(s, Inches(9.1), Inches(3.05), Inches(3.5), Inches(3.0),
             money_quote, font=FONT_BODY, size=13, italic=True,
             color=INK, line_spacing=1.45)
    add_footer(s)
    add_right_meta(s)
    return s


def build_ci_pipeline(prs):
    s = build_content(
        prs,
        "Section 03  ·  The CI pipeline",
        "Documents are one half. Tests are the other.",
        [
            {"x": Inches(0.55), "y": Inches(2.2), "w": Inches(12), "h": Inches(0.5),
             "text": "On every push to GitHub:", "size": 14, "bold": True, "color": MUTED},
            {"x": Inches(0.55), "y": Inches(2.7), "w": Inches(12), "h": Inches(3.5),
             "text": ("1.  Parse check — catches the duplicate-let class of bug.\n"
                      "2.  ESLint — coding-hygiene gate.\n"
                      "3.  Static secret scan — fails the build if any deployed file leaks credentials.\n"
                      "4.  33 unit tests on the computation module — toMmol, TIR, trend, smoothing.\n"
                      "5.  npm audit — refreshes the SOUP CVE table for §8 of the cybersecurity doc.\n"
                      "6.  SBOM regenerated and published as a workflow artifact."),
             "size": 15, "line_spacing": 1.55},
            {"x": Inches(0.55), "y": Inches(6.4), "w": Inches(12), "h": Inches(0.5),
             "text": "End-to-end run time: ~2 minutes. Result is the V&V evidence trail.",
             "size": 13, "italic": True, "color": MUTED},
        ],
    )
    return s


def build_ai_vs_human(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    add_top_divider(s, "Section 04  ·  Where AI ends, humans begin")
    add_text(s, Inches(0.55), Inches(1.0), Inches(12), Inches(1),
             "AI compresses drafting. Humans own accountability.",
             font=FONT_HEAD, size=32, bold=True, color=INK,
             line_spacing=1.0)
    # left column
    add_rect(s, Inches(0.55), Inches(2.4), Inches(6.0), Inches(4.4), WHITE)
    add_text(s, Inches(0.85), Inches(2.55), Inches(5.5), Inches(0.4),
             "AI does well", font=FONT_HEAD, size=14, bold=True, color=PURPLE_100)
    add_text(s, Inches(0.85), Inches(3.0), Inches(5.5), Inches(3.7),
             ("•  Document scaffolding from standards.\n\n"
              "•  Cross-referencing — REQ / DES / RISK / TST IDs stay consistent across nine docs.\n\n"
              "•  Code-grounded honesty — flags real bugs, real gaps as Open Items.\n\n"
              "•  Drafting velocity — order of magnitude faster than manual authorship."),
             font=FONT_BODY, size=13, color=INK, line_spacing=1.4)
    # right column
    add_rect(s, Inches(6.75), Inches(2.4), Inches(6.0), Inches(4.4), WHITE)
    add_text(s, Inches(7.05), Inches(2.55), Inches(5.5), Inches(0.4),
             "AI does NOT replace", font=FONT_HEAD, size=14, bold=True, color=PEACH_100)
    add_text(s, Inches(7.05), Inches(3.0), Inches(5.5), Inches(3.7),
             ("•  Clinical severity calibration.\n\n"
              "•  Use-error analysis (IEC 62366).\n\n"
              "•  QMS integration — your templates, numbering, approval workflow.\n\n"
              "•  Sign-off and accountability — real signatures.\n\n"
              "•  Penetration testing.\n\n"
              "•  Formal SOUP qualification."),
             font=FONT_BODY, size=13, color=INK, line_spacing=1.4)
    add_footer(s)
    add_right_meta(s)
    return s


def build_quantified_value(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, GREY_100)
    add_top_divider(s, "Section 05  ·  Quantified value")
    add_text(s, Inches(0.55), Inches(1.0), Inches(12), Inches(1),
             "Order-of-magnitude compression",
             font=FONT_HEAD, size=36, bold=True, color=INK)
    # Two big numbers
    add_text(s, Inches(0.55), Inches(2.5), Inches(6.2), Inches(0.5),
             "AI generation time", font=FONT_HEAD, size=14, bold=True, color=MUTED)
    add_text(s, Inches(0.55), Inches(3.0), Inches(6.2), Inches(2.0),
             "~30 min", font=FONT_HEAD, size=96, bold=True, color=PURPLE_100,
             line_spacing=0.95)
    add_text(s, Inches(0.55), Inches(5.2), Inches(6.2), Inches(0.6),
             "across the nine documents",
             font=FONT_BODY, size=13, italic=True, color=MUTED)

    add_text(s, Inches(7.0), Inches(2.5), Inches(6.0), Inches(0.5),
             "Equivalent manual effort", font=FONT_HEAD, size=14, bold=True, color=MUTED)
    add_text(s, Inches(7.0), Inches(3.0), Inches(6.0), Inches(2.0),
             "25–45 d", font=FONT_HEAD, size=96, bold=True, color=PEACH_100,
             line_spacing=0.95)
    add_text(s, Inches(7.0), Inches(5.2), Inches(6.0), Inches(0.6),
             "person-days, plus clinical and security review",
             font=FONT_BODY, size=13, italic=True, color=MUTED)

    add_text(s, Inches(0.55), Inches(6.3), Inches(12), Inches(0.5),
             "Plus freshness: every code change re-runs the pipeline. The DHF stays alive instead of frozen at submission.",
             font=FONT_BODY, size=14, italic=True, color=INK)
    add_footer(s)
    add_right_meta(s)
    return s


def build_qa(prs):
    s = build_content(
        prs,
        "Section 06  ·  Q&A",
        "Three asks, then your questions.",
        [
            {"x": Inches(0.55), "y": Inches(2.4), "w": Inches(12), "h": Inches(0.6),
             "text": "1. Validate quality.", "size": 18, "bold": True, "color": PURPLE_100},
            {"x": Inches(0.85), "y": Inches(2.95), "w": Inches(12), "h": Inches(0.8),
             "text": "Pick one document — Risk Analysis is usually the most stressful — and have RA red-team it.",
             "size": 14, "color": INK},

            {"x": Inches(0.55), "y": Inches(4.0), "w": Inches(12), "h": Inches(0.6),
             "text": "2. Test the boundaries.", "size": 18, "bold": True, "color": PURPLE_100},
            {"x": Inches(0.85), "y": Inches(4.55), "w": Inches(12), "h": Inches(0.8),
             "text": "Pick a Trinity product currently in development. Can AI scaffold its DHF the same way? Run the experiment.",
             "size": 14, "color": INK},

            {"x": Inches(0.55), "y": Inches(5.6), "w": Inches(12), "h": Inches(0.6),
             "text": "3. Decide on QMS integration.", "size": 18, "bold": True, "color": PURPLE_100},
            {"x": Inches(0.85), "y": Inches(6.15), "w": Inches(12), "h": Inches(0.8),
             "text": "Templates, numbering schemes, approval flows mapped to your QMS. 1–2 week engagement.",
             "size": 14, "color": INK},
        ],
    )
    return s


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title(prs)
    build_agenda(prs)

    build_section_divider(prs, "01", ["The", "Product"])
    build_product_context(prs)

    build_section_divider(prs, "02", ["DHF", "Artifact Set"])
    build_dhf_overview(prs)

    build_doc_card(
        prs, "01", "Software Safety Classification",
        "IEC 62304 §4.3  ·  TCG-REG-01",
        [
            "Class B for the demo, with explicit rationale.",
            "Class C upgrade path documented for any treatment-decision claim.",
            "External controls (controlled distribution, no clinical labeling) reduce demo residual risk.",
            "Open Items flag four judgment calls for human RA.",
        ],
        "AI flags the question — \"if you ever claim treatment-decision support, redo this DHF under Class C\" — rather than burying it.",
    )

    build_doc_card(
        prs, "02", "Software Requirements Specification",
        "IEC 62304 §5.2  ·  TCG-REG-02",
        [
            "67 requirements with REQ-IDs across functional, performance, interface, safety, security, usability.",
            "REQ-080: \"glucose conversion shall use divisor 18.0182 with no rounding\" — extracted from actual code.",
            "REQ-094: \"no console output containing the access token\" — from AI's own threat model.",
            "Each REQ has a TST-NNN verification reference. The tests are real and run in CI.",
        ],
        "Not boilerplate. Each requirement is grounded in code or a real risk the AI identified.",
    )

    build_doc_card(
        prs, "03", "Software Design Description",
        "IEC 62304 §5.3  ·  TCG-REG-03",
        [
            "Architectural decomposition: 7 client items, 2 server items.",
            "Trust-boundary diagram. Data dictionary. SOUP component table.",
            "Open Items flags six existing bugs that should drive change orders.",
            "Calls out code-quality issues your software lead would catch — and writes them down.",
        ],
        "AI is reading code and writing the list of code-quality issues a senior engineer would write.",
    )

    build_doc_card(
        prs, "04", "Risk Analysis  (centerpiece)",
        "ISO 14971:2019  ·  TCG-REG-04",
        [
            "18 identified risks with severity / probability / risk-class scoring.",
            "RISK-04: credentials in client — initial HIGH, residual LOW post-mitigations we shipped.",
            "RISK-07: TIR computed row-weighted not time-weighted — a real bug, found by AI in our actual code.",
            "Distinguishes implemented RCMs from PROPOSED — Future. Honest about gaps.",
            "Severity scaling needs clinical input. Probability needs field data. AI flags this.",
        ],
        "Honest about residual risk. Tells us where it can't help. The framing your RA team would write themselves.",
    )

    build_doc_card(
        prs, "05", "Software Test Description",
        "IEC 62304 §5.5–5.7  ·  TCG-REG-05",
        [
            "35 test cases: parse, lint, secret-scan, 10 unit, 9 integration, 12 system, 4 manual.",
            "Each test traces to one or more REQ-IDs and corresponding RISK-IDs.",
            "Unit tests on the computation module are the high-leverage Class B verification activity.",
            "TST-UNIT-001 verifies toMmol(18.0182) returns exactly 1.0 within IEEE 754 precision.",
        ],
        "These are real tests. They actually run in our CI on every push.",
    )

    build_doc_card(
        prs, "06", "Cybersecurity Assessment",
        "FDA Premarket Cyber 2023  ·  TCG-REG-06",
        [
            "STRIDE per trust boundary. Asset inventory. Security control catalog.",
            "SOUP CVE status snapshot, with explicit pointer to npm audit as live source of truth.",
            "Postmarket update plan placeholder, with what productization would require.",
            "Honest scope statement on what relies on Netlify, BioT, browser/OS vendor security programs.",
        ],
        "AI was honest that the SOUP CVE list is its static snapshot. It's not pretending to be a vulnerability scanner.",
    )

    build_doc_card(
        prs, "07", "Software Bill of Materials",
        "CycloneDX 1.5  ·  TCG-REG-07",
        [
            "Real components, real versions, real hashes from package-lock.json.",
            "External services (BioT, Netlify, WebAuthn platform) included as service components.",
            "Mechanically regenerated by CI on every push. Always current.",
            "Required by FDA 2023 cybersecurity guidance.",
        ],
        "Mechanical artifact done right. The kind of thing AI should fully own.",
    )

    build_doc_card(
        prs, "08", "Traceability Matrix",
        "IEC 62304 implicit  ·  TCG-REG-08",
        [
            "Bidirectional: REQ → DES → TEST → RISK and reverse.",
            "Coverage summary: 56 / 67 requirements have automated test coverage.",
            "11 explicit gaps listed with recommended additions.",
            "RISK → CONTROL → TEST tracks every risk to its mitigation evidence.",
        ],
        "The single biggest pain point of any regulatory team. AI keeps it consistent because it generated all source documents in one session.",
    )

    build_section_divider(prs, "03", ["The", "CI Pipeline"])
    build_ci_pipeline(prs)

    build_section_divider(prs, "04", ["Where AI Ends,", "Humans Begin"])
    build_ai_vs_human(prs)

    build_section_divider(prs, "05", ["Quantified", "Value"])
    build_quantified_value(prs)

    build_section_divider(prs, "06", ["Q & A"])
    build_qa(prs)

    build_thank_you(prs)

    out = "presentations/trinity-cgm-plus-ai-regulatory.pptx"
    import os
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)
    print(f"Wrote {out} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
